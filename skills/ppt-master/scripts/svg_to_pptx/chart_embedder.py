"""Embed native PowerPoint charts extracted from SVG chart pages."""

from __future__ import annotations

import tempfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu

from .pptx_dimensions import EMU_PER_PIXEL

try:
    from svg_to_excel_chart import (
        extract_chart_data,
        is_data_heavy_chart,
        list_charts_in_svg,
        svg_chart_to_excel,
    )
except ImportError:
    extract_chart_data = None  # type: ignore[misc, assignment]
    is_data_heavy_chart = None  # type: ignore[misc, assignment]
    list_charts_in_svg = None  # type: ignore[misc, assignment]
    svg_chart_to_excel = None  # type: ignore[misc, assignment]


def _plot_area_box(
    summary: dict[str, Any],
    pixel_width: int,
    pixel_height: int,
) -> tuple[int, int, int, int]:
    plot = summary.get('plot_area') or {}
    if plot.get('kind') == 'rect':
        x_min = float(plot.get('x_min', pixel_width * 0.1))
        y_min = float(plot.get('y_min', pixel_height * 0.15))
        x_max = float(plot.get('x_max', pixel_width * 0.9))
        y_max = float(plot.get('y_max', pixel_height * 0.85))
    else:
        cx = float(plot.get('center_x') or pixel_width * 0.35)
        cy = float(plot.get('center_y') or pixel_height * 0.5)
        radius = float(plot.get('radius') or min(pixel_width, pixel_height) * 0.2)
        x_min = cx - radius
        x_max = cx + radius
        y_min = cy - radius
        y_max = cy + radius

    left = int(x_min * EMU_PER_PIXEL)
    top = int(y_min * EMU_PER_PIXEL)
    width = int(max(x_max - x_min, 1) * EMU_PER_PIXEL)
    height = int(max(y_max - y_min, 1) * EMU_PER_PIXEL)
    return left, top, width, height


def _chart_type_enum(chart_type: str) -> XL_CHART_TYPE:
    if chart_type == 'line':
        return XL_CHART_TYPE.LINE
    if chart_type == 'pie':
        return XL_CHART_TYPE.PIE
    return XL_CHART_TYPE.COLUMN_CLUSTERED


def _should_embed_chart(summary: dict[str, Any], chart_mode: str) -> bool:
    if chart_mode == 'svg':
        return False
    if not summary.get('extractable'):
        return False
    if chart_mode == 'excel':
        return True
    if chart_mode == 'hybrid' and is_data_heavy_chart is not None:
        return is_data_heavy_chart(summary)
    return False


def _build_chart_data(summary: dict[str, Any]) -> ChartData:
    chart_data = ChartData()
    categories = summary.get('categories') or []
    chart_data.categories = categories
    series_list = summary.get('series') or []
    if not series_list:
        chart_data.add_series('Series 1', tuple([0.0] * len(categories)))
        return chart_data
    for series in series_list:
        values = series.get('values') or []
        padded = list(values) + [0.0] * max(0, len(categories) - len(values))
        chart_data.add_series(series.get('name') or 'Series 1', tuple(padded[:len(categories)]))
    return chart_data


def embed_charts_in_pptx(
    pptx_path: Path,
    svg_files: list[Path],
    *,
    chart_mode: str = 'svg',
    pixel_width: int,
    pixel_height: int,
    verbose: bool = True,
    temp_dir: Path | None = None,
) -> dict[str, Any]:
    """Post-process a PPTX, replacing extractable SVG charts with native chart objects."""
    report: dict[str, Any] = {
        'chart_mode': chart_mode,
        'slides': [],
        'embedded_count': 0,
        'skipped_count': 0,
    }

    if chart_mode == 'svg' or extract_chart_data is None:
        return report

    if not pptx_path.exists():
        raise FileNotFoundError(f'PPTX not found: {pptx_path}')

    work_temp = temp_dir or Path(tempfile.mkdtemp())
    xlsx_dir = work_temp / 'chart_xlsx'
    xlsx_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(pptx_path))
    if len(prs.slides) != len(svg_files):
        if verbose:
            print(
                f'  [chart] Slide count mismatch ({len(prs.slides)} vs {len(svg_files)}); '
                'skipping chart embedding',
            )
        return report

    for slide_idx, svg_path in enumerate(svg_files):
        slide_report: dict[str, Any] = {
            'svg': svg_path.name,
            'charts': [],
        }
        try:
            summaries = list_charts_in_svg(svg_path) if list_charts_in_svg else []
        except Exception as exc:
            slide_report['error'] = str(exc)
            report['slides'].append(slide_report)
            continue

        if not summaries:
            report['slides'].append(slide_report)
            continue

        slide = prs.slides[slide_idx]
        for summary in summaries:
            embed = _should_embed_chart(summary, chart_mode)
            chart_entry = {
                'chart_type': summary.get('chart_type'),
                'chart_element_id': summary.get('chart_element_id'),
                'embedded': embed,
                'data_point_count': summary.get('data_point_count'),
            }
            if not embed:
                report['skipped_count'] += 1
                slide_report['charts'].append(chart_entry)
                continue

            xlsx_path = xlsx_dir / f'{svg_path.stem}_chart{summary.get("chart_index", 0)}.xlsx'
            try:
                if svg_chart_to_excel is not None:
                    svg_chart_to_excel(
                        svg_path,
                        xlsx_path,
                        chart_type=str(summary.get('chart_type') or 'auto'),
                        chart_index=int(summary.get('chart_index') or 0),
                    )
                    chart_entry['xlsx'] = str(xlsx_path)
            except Exception as exc:
                chart_entry['error'] = str(exc)
                report['skipped_count'] += 1
                slide_report['charts'].append(chart_entry)
                continue

            left, top, width, height = _plot_area_box(summary, pixel_width, pixel_height)
            chart_data = _build_chart_data(summary)
            chart_type = str(summary.get('chart_type') or 'bar')
            slide.shapes.add_chart(
                _chart_type_enum(chart_type),
                left, top, width, height,
                chart_data,
            )
            report['embedded_count'] += 1
            slide_report['charts'].append(chart_entry)

        report['slides'].append(slide_report)

    prs.save(str(pptx_path))

    if verbose and (report['embedded_count'] or report['skipped_count']):
        print(
            f'  Chart mode ({chart_mode}): embedded {report["embedded_count"]}, '
            f'skipped {report["skipped_count"]}',
        )

    return report
